#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Mon Oct 12 15:54:20 2020

@author: ssccdmx-amh
"""
from pptx import Presentation
# para imagenegenes se utiliza este modulo
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

class ArchivosPptx():    
    def titulo(self,prst,titulo):
        prs = prst
        tituDiapo = prs.slide_layouts[0]       
        slide = prs.slides.add_slide(tituDiapo)        
        title = slide.shapes.title             
        title.text = titulo
         
        
    def segundoFormato(self, prst):
        prs = prst
        bullet_slide_layout = prs.slide_layouts[1]
        
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = 'Diapositiva de viñeta'
        
        tf = body_shape.text_frame
        tf.text = 'Find the bullet slide layout'
        
        p = tf.add_paragraph()
        p.text = 'Use _TextFrame.text for first bullet'
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
        p.level = 2
        
    def agregarImagen(self,templete,pathImg):
        # Ruta de imagen
        img_path = pathImg

        prs = Presentation(templete)
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        left = top = Inches(1)
        pic = slide.shapes.add_picture(img_path, left, top)
        
        left = Inches(5)
        height = Inches(5.5)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
        prs.save('Salida.pptx')
        
    def crearTabla(self):
        
        from pptx.util import Inches
        prs = Presentation()
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes
        
        shapes.title.text = 'Añadiendo tablas'
        
        rows = cols = 2
        left = top = Inches(2.0)
        width = Inches(6.0)
        height = Inches(0.8)
        
        table = shapes.add_table(rows, cols, left, top, width, height).table
        
        # set column widths
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(4.0)
        
        # write column headings
        table.cell(0, 0).text = 'Encabezado-1'
        table.cell(0, 1).text = 'Encabezado-2'
        
        # write body cells
        table.cell(1, 0).text = 'valo-1'
        table.cell(1, 1).text = 'valor-2'
        
        prs.save('creacionTabla.pptx')
        
    def diseñoDiapo1(self,prs):        
        SLD_LAYOUT_TITLE_AND_CONTENT = 1        
        slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title        
        title.text = "Diapositiva"        
        subtitle = slide.placeholders[1]
        subtitle.text = "subtitulo!"
        
    def diseñoDiapo2(self,prs):
        # =====================================
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)        
        left = top = width = height = Inches(1)
        
        title = slide.shapes.title        
        title.text = "Diapositiva"
        
        txBox = slide.shapes.add_textbox(Inches(4.07), Inches(1.93), Inches(1.82), Inches(5.54))        
        tf = txBox.text_frame                
        tf.text = "Esto es un parrafo de texto \n"+"Hola mundo"        
        p = tf.add_paragraph()
        # p.text = "This is a second paragraph that's bold"
        # p.font.bold = True
        # p = tf.add_paragraph()
        # p.text = "This is a third paragraph that's big"
        # p.font.size = Pt(30)
        #========================================
        img_path = 'descarga.png'             
        left = Inches(1.00)
        top = Inches(1.67)
        height = Inches(2.45)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
        # left = Inches(1)
        # height = Inches(1)
        # pic = slide.shapes.add_picture(img_path, left, top, height=height) 
        # left = top = Inches(1)
        # pic = slide.shapes.add_picture(img_path, left, top)
        
            
        
    def abrirPlantilla(self,ruta):
        f = open(ruta)
        prs = Presentation(ruta)        
        self.titulo(prs,'Presentacion')                
                
        self.diseñoDiapo1(prs)
        self.diseñoDiapo2(prs)
        
        
# =============================================================================
#         shapes = slide.shapes
#         left = top = width = height = Inches(1.0)
#         shape = shapes.add_shape(
#             MSO_SHAPE.CUBE, left, top, width, height
#         )
# =============================================================================
        
# =============================================================================
#         shape.text = 'foobar'
# 
#         # is equivalent to ...
#         
#         text_frame = shape.text_frame
#         text_frame.clear()
#         p = text_frame.paragraphs[0]
#         run = p.add_run()
#         run.text = 'foobar'
# =============================================================================

        
        
                
        prs.save('Salida.pptx')    
                
        f.close()

        

miObjeto = ArchivosPptx()
miObjeto.abrirPlantilla('templete.pptx')